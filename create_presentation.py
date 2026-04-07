from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
COLOR_PRIMARY = RGBColor(26, 31, 46)  # Dark blue
COLOR_ACCENT = RGBColor(102, 126, 234)  # Purple/Blue
COLOR_SECONDARY = RGBColor(76, 175, 80)  # Green
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LIGHT_GRAY = RGBColor(224, 224, 224)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = COLOR_ACCENT
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add decorative line
    line = slide.shapes.add_shape(1, Inches(2), Inches(4), Inches(6), Inches(0))
    line.line.color.rgb = COLOR_ACCENT
    line.line.width = Pt(3)

def add_content_slide(prs, title, content_items, use_colors=False):
    """Add a slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Add header background bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.color.rgb = COLOR_PRIMARY
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_WHITE
    
    # Add accent line
    accent_line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(3), Inches(0))
    accent_line.line.color.rgb = COLOR_ACCENT
    accent_line.line.width = Pt(4)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    colors = [COLOR_ACCENT, COLOR_SECONDARY, RGBColor(255, 152, 0)]
    
    for idx, item in enumerate(content_items):
        if idx > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[idx]
        p.text = "• " + item
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_PRIMARY
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        
        if use_colors:
            p.font.color.rgb = colors[idx % len(colors)]

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Add header
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_WHITE
    
    accent_line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(3), Inches(0))
    accent_line.line.color.rgb = COLOR_ACCENT
    accent_line.line.width = Pt(4)
    
    # Left column
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4), Inches(0.4))
    left_title_frame = left_title_box.text_frame
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.text = left_title
    left_title_para.font.size = Pt(18)
    left_title_para.font.bold = True
    left_title_para.font.color.rgb = COLOR_ACCENT
    
    left_content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.2), Inches(5))
    left_frame = left_content_box.text_frame
    left_frame.word_wrap = True
    
    for idx, item in enumerate(left_items):
        if idx > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[idx]
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_PRIMARY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Right column
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.4), Inches(4), Inches(0.4))
    right_title_frame = right_title_box.text_frame
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.text = right_title
    right_title_para.font.size = Pt(18)
    right_title_para.font.bold = True
    right_title_para.font.color.rgb = COLOR_SECONDARY
    
    right_content_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4.2), Inches(5))
    right_frame = right_content_box.text_frame
    right_frame.word_wrap = True
    
    for idx, item in enumerate(right_items):
        if idx > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[idx]
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_PRIMARY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Add divider line
    divider = slide.shapes.add_shape(1, Inches(5), Inches(1.4), Inches(0), Inches(5.2))
    divider.line.color.rgb = COLOR_LIGHT_GRAY
    divider.line.width = Pt(2)

# Slide 1: Title Slide
add_title_slide(prs, "FitPulse", "Detecting Unusual Health Patterns Using Fitness Watch Data")

# Slide 2: Project Overview
add_content_slide(prs, "Project Overview", [
    "Analyzes wearable fitness device data to detect anomalies in health patterns",
    "Leverages AI-based anomaly detection for intelligent health monitoring",
    "Enables proactive health monitoring and preventive healthcare",
    "Supports personalized wellness insights and healthcare collaboration"
], use_colors=True)

# Slide 3: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "Surge in wearable fitness devices generates vast time-series health data",
    "Users and healthcare providers often miss subtle early warning signs",
    "Anomalies like irregular heartbeats, sleep disturbances go unnoticed",
    "Need for intelligent analysis to flag unusual health patterns automatically"
], use_colors=True)

# Slide 4: Key Outcomes
add_content_slide(prs, "Project Outcomes", [
    "Accurate detection of anomalies in heart rate, sleep, and step count",
    "Personalized health alerts using time-series models and clustering",
    "Integration with raw fitness watch data (CSV/JSON format)",
    "Real-time or batch-based anomaly flagging for dynamic use cases",
    "Interactive dashboards for visual trend and anomaly tracking",
    "Exportable reports for users and healthcare professionals"
], use_colors=True)

# Slide 5: Module 1 - Data Collection & Preprocessing
add_content_slide(prs, "Module 1: Data Collection & Preprocessing", [
    "Import health data (heart rate, steps, sleep) from fitness trackers",
    "Support CSV/JSON format ingestion logic",
    "Clean and normalize timestamps (UTC conversion)",
    "Handle missing/null values through interpolation",
    "Align time intervals to consistent frequency (1-minute granularity)"
], use_colors=True)

# Slide 6: Module 2 - Feature Extraction & Modeling
add_content_slide(prs, "Module 2: Feature Extraction & Modeling", [
    "Extract statistical features using TSFresh (mean, std, kurtosis)",
    "Apply Facebook Prophet for trend modeling and deviation detection",
    "Use clustering (KMeans, DBSCAN) for behavioral pattern detection",
    "Identify seasonal trends and anomalous deviations",
    "Generate feature matrices for anomaly detection"
], use_colors=True)

# Slide 7: Module 3 - Anomaly Detection & Visualization
add_content_slide(prs, "Module 3: Anomaly Detection & Visualization", [
    "Rule-based anomaly detection via threshold violations",
    "Model-based detection using Prophet residuals",
    "Clustering-based outlier identification",
    "Interactive visualization with Matplotlib/Plotly",
    "Annotated charts highlighting anomalies with time windows"
], use_colors=True)

# Slide 8: Module 4 - Dashboard for Insights
add_content_slide(prs, "Module 4: Dashboard for Insights", [
    "Build interactive Streamlit-based user interface",
    "Enable dynamic file upload and anomaly detection triggering",
    "Filter data by date, metric, and health parameters",
    "Generate comprehensive anomaly summaries and trend reports",
    "Export insights in PDF/CSV format for healthcare use"
], use_colors=True)

# Slide 9: Milestone 1 - Weeks 1-2
add_two_column_slide(prs, "Milestone 1: Data Collection & Preprocessing (Weeks 1-2)",
    "Requirements",
    [
        "CSV/JSON ingestion logic",
        "Timestamp normalization",
        "Missing value handling",
        "Time-aligned resampling"
    ],
    "Deliverables",
    [
        "File upload UI",
        "Cleaned dataset preview",
        "Time-normalized data log",
        "Data quality report"
    ]
)

# Slide 10: Milestone 2 - Weeks 3-4
add_two_column_slide(prs, "Milestone 2: Feature Extraction & Modeling (Weeks 3-4)",
    "Requirements",
    [
        "TSFresh feature extraction",
        "Prophet trend modeling",
        "Clustering algorithms",
        "Feature matrix generation"
    ],
    "Deliverables",
    [
        "TSFresh feature matrix",
        "Prophet trend graphs",
        "Clustering visualizations",
        "PCA/t-SNE projections"
    ]
)

# Slide 11: Milestone 3 - Weeks 5-6
add_two_column_slide(prs, "Milestone 3: Anomaly Detection & Visualization (Weeks 5-6)",
    "Requirements",
    [
        "Multi-method detection",
        "Threshold rules engine",
        "Outlier identification",
        "Interactive charts"
    ],
    "Deliverables",
    [
        "Heart rate anomaly charts",
        "Sleep pattern visualization",
        "Step count alerts",
        "Annotated visualizations"
    ]
)

# Slide 12: Milestone 4 - Weeks 7-8
add_two_column_slide(prs, "Milestone 4: Dashboard & Insights (Weeks 7-8)",
    "Requirements",
    [
        "Streamlit interface",
        "File upload mechanism",
        "Dynamic processing",
        "Report generation"
    ],
    "Deliverables",
    [
        "Interactive dashboard UI",
        "Real-time anomaly alerts",
        "Downloadable reports",
        "Date/metric filters"
    ]
)

# Slide 13: Technology Stack
add_two_column_slide(prs, "Technology Stack",
    "Data & ML Libraries",
    [
        "Pandas - Data manipulation",
        "NumPy - Numerical computing",
        "TSFresh - Time series features",
        "Prophet - Trend forecasting",
        "Scikit-learn - Clustering & ML"
    ],
    "Visualization & UI",
    [
        "Streamlit - Interactive dashboard",
        "Matplotlib/Seaborn - Static plots",
        "Plotly - Interactive charts",
        "Python 3.x environment",
        "CSV/JSON data formats"
    ]
)

# Slide 14: Architecture Overview
add_content_slide(prs, "System Architecture", [
    "Data Layer: CSV/JSON import and preprocessing engine",
    "Feature Layer: TSFresh extraction and Prophet modeling",
    "Detection Layer: Rule-based and ML-based anomaly detection",
    "Visualization Layer: Plotly/Matplotlib for interactive charts",
    "Dashboard Layer: Streamlit UI for user interaction and reporting"
], use_colors=True)

# Slide 15: Evaluation Criteria
add_two_column_slide(prs, "Evaluation Criteria",
    "Data Quality & Features",
    [
        "M1: Successful data import & cleaning",
        "M1: Consistent time-series format",
        "M2: TSFresh features extracted",
        "M2: Prophet trends visible"
    ],
    "Detection & Reporting",
    [
        "M3: 90%+ anomaly detection accuracy",
        "M3: Visual alignment with patterns",
        "M4: Dashboard functions correctly",
        "M4: Report generation works"
    ]
)

# Slide 16: Key Features
add_content_slide(prs, "Key Features", [
    "Multi-metric anomaly detection (heart rate, sleep, steps)",
    "Configurable thresholds for personalized alerts",
    "Real-time and batch processing capabilities",
    "Historical trend analysis and pattern recognition",
    "User-friendly interactive dashboard interface",
    "Exportable reports for healthcare professionals"
], use_colors=True)

# Slide 17: Expected Impact
add_content_slide(prs, "Expected Impact", [
    "Early detection of health anomalies enables preventive healthcare",
    "Personalized alerts help users make informed wellness decisions",
    "Healthcare providers gain actionable insights from fitness data",
    "Scalable solution for population health monitoring",
    "Integration potential with healthcare systems and wearables",
    "Improved quality of life through proactive health management"
], use_colors=True)

# Slide 18: Timeline
add_two_column_slide(prs, "Project Timeline",
    "Phase 1 (Weeks 1-2)",
    [
        "Module 1: Data Collection",
        "Data ingestion setup",
        "Preprocessing pipeline"
    ],
    "Phase 2 (Weeks 3-4)",
    [
        "Module 2: Feature Extraction",
        "TSFresh modeling",
        "Clustering implementation"
    ]
)

# Slide 19: Timeline (continued)
add_two_column_slide(prs, "Project Timeline (Continued)",
    "Phase 3 (Weeks 5-6)",
    [
        "Module 3: Anomaly Detection",
        "Multi-method detection",
        "Visualization development"
    ],
    "Phase 4 (Weeks 7-8)",
    [
        "Module 4: Dashboard",
        "Streamlit integration",
        "Final testing & deployment"
    ]
)

# Slide 20: Closing Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_PRIMARY

# Add centered message
text_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p1 = text_frame.paragraphs[0]
p1.text = "FitPulse"
p1.font.size = Pt(56)
p1.font.bold = True
p1.font.color.rgb = COLOR_ACCENT
p1.alignment = PP_ALIGN.CENTER

text_frame.add_paragraph()
p2 = text_frame.paragraphs[1]
p2.text = "Proactive Health Monitoring Through Intelligence"
p2.font.size = Pt(28)
p2.font.color.rgb = COLOR_WHITE
p2.alignment = PP_ALIGN.CENTER

text_frame.add_paragraph()
p3 = text_frame.paragraphs[2]
p3.text = "\n"
p3.font.size = Pt(12)

text_frame.add_paragraph()
p4 = text_frame.paragraphs[3]
p4.text = "Thank You"
p4.font.size = Pt(32)
p4.font.bold = True
p4.font.color.rgb = COLOR_SECONDARY
p4.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = os.path.join(os.path.dirname(__file__), "FitPulse_Project_Presentation.pptx")
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")
